"""Gera o PDF do Fechamento do Mercado com os dados atualizados"""

import datetime
from typing import Literal
from pathlib import Path
import locale

from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from settings import BASE_DIR

def cria_pdf(valores_dos_ativos: dict, acoes_em_alta_baixa: dict, modelo: Path, salvar: Path) -> None:
    """Gera o PDF do Fechamento de Mercado com os dados atualizados

    ### Argumentos:
    - valores_dos_ativos: Dicionário contendo os dados dos ativos
    - acoes_em_alta_baixa: Dicionário contendo as ações em alta e baixa"""

    apresentacao = Presentation(modelo)
    slide = apresentacao.slides[0]

    locale.setlocale(locale.LC_TIME, 'pt_BR.UTF-8')

    for shape in slide.shapes:

        match shape.name:
            case 'Variação IBOV':
                formatacao_textos_pp(shape, valores_dos_ativos["ativos"][0]["Variação"], 327, 279, "numérico", True)
                inserir_negativo_positivo(slide, valores_dos_ativos["ativos"][0]["Variação"], 285, 287, 23, 22)

            case 'Cotação IBOVESPA':
                formatacao_textos_pp(shape, round(valores_dos_ativos["ativos"][0]["Valor"]), 333, 306, "ibov", False)

            case 'Variação NASDAQ':
                formatacao_textos_pp(shape, valores_dos_ativos["ativos"][4]["Variação"], 327, 343, "numérico", True)
                inserir_negativo_positivo(slide, valores_dos_ativos["ativos"][4]["Variação"], 285, 351, 23, 22)

            case 'Variação S&P 500':
                formatacao_textos_pp(shape, valores_dos_ativos["ativos"][1]["Variação"], 327, 396, "numérico", True)
                inserir_negativo_positivo(slide, valores_dos_ativos["ativos"][1]["Variação"], 285, 402, 23, 22)

            case 'Variação EUROSTOXX':
                formatacao_textos_pp(shape, valores_dos_ativos["ativos"][2]["Variação"], 327, 454, "numérico", True)
                inserir_negativo_positivo(slide, valores_dos_ativos["ativos"][2]["Variação"], 285, 461, 23, 22)

            case 'Cotação DÓLAR':
                formatacao_textos_pp(shape, valores_dos_ativos["ativos"][3]["Valor"], 327, 516, "dolar", True)
                inserir_negativo_positivo(slide, valores_dos_ativos["ativos"][3]["Valor"], 285, 521, 23, 22)

            case 'Variação maior alta 1':
                formatacao_textos_pp(shape, acoes_em_alta_baixa["alta"][0]["Variação"], 327, 630, "numérico", True)
                inserir_negativo_positivo(slide, acoes_em_alta_baixa["alta"][0]["Variação"], 285, 636, 23, 22)

            case 'Variação maior alta 2':
                formatacao_textos_pp(shape, acoes_em_alta_baixa["alta"][1]["Variação"], 327, 677, "numérico", True)
                inserir_negativo_positivo(slide, acoes_em_alta_baixa["alta"][1]["Variação"], 285, 684, 23, 22)

            case 'Variação maior queda 1':
                formatacao_textos_pp(shape, acoes_em_alta_baixa["baixa"][0]["Variação"], 327, 796, "numérico", True)
                inserir_negativo_positivo(slide, acoes_em_alta_baixa["baixa"][0]["Variação"], 285, 803, 23, 22)

            case 'Variação maior queda 2':
                formatacao_textos_pp(shape, acoes_em_alta_baixa["baixa"][1]["Variação"], 327, 842,  "numérico", True)
                inserir_negativo_positivo(slide, acoes_em_alta_baixa["baixa"][1]["Variação"], 285, 849, 23, 22)

            case 'Ativo maior alta 1':
                formatacao_textos_pp(shape, acoes_em_alta_baixa["alta"][0]["Ticker"], 119, 630, "ticker", True)

            case 'Ativo maior alta 2':
                formatacao_textos_pp(shape, acoes_em_alta_baixa["alta"][1]["Ticker"], 119, 677, "ticker", True)

            case 'Ativo maior queda 1':
                formatacao_textos_pp(shape, acoes_em_alta_baixa["baixa"][0]["Ticker"], 119, 796, "ticker", True)

            case 'Ativo maior queda 2':
                formatacao_textos_pp(shape, acoes_em_alta_baixa["baixa"][1]["Ticker"], 119, 842, "ticker", True)

            case 'Data':
                formatacao_textos_pp(shape, datetime.datetime.today().strftime('%d de %B de %Y'), 173, 220, "data")

    apresentacao.save(salvar)

def formatacao_textos_pp(
    shape: SlideShapes,
    texto: float | int | str,
    posicao_left: Pt,
    posicao_top: Pt,
    formatacao: Literal[
        "ibov",
        "data",
        "dolar",
        "ticker",
        "numérico",
    ],
    bold: bool = True,
) -> None:
    """Formata o dataframe dos ativos em alta e em baixa e transforma em dicionário.

    ### Argumentos:
    - shape (SlidesShape): Indentifica as formas no slide do pdf
    - texto(float | int | str): Indetifica se é o valor do ativo ou o ticker
    - posicao_left(Pt): Posição horizontal da forma no slide
    - posicao_top(Pt): Posição vertical da forma no slide
    - formatacao(Literal): Indentifica a formatação que vai ser usada
    - bold (bool): Indica se o texto é em negrito"""

    paragrafo = shape.text_frame.paragraphs[0]

    # ? - Formatando ticker
    if isinstance(texto, float | int):
        texto = (
            str(round(texto, 2))
            .replace(".", ",")
            .replace("-", "")
        )

    match formatacao:
        case "numérico":
            tamanho_fonte = Pt(22)
            nome_fonte = 'Segoe UI'
            cor_fonte = RGBColor(115, 111, 111)
            texto = texto + "%"

        case "ticker":
            tamanho_fonte = Pt(22)
            nome_fonte = 'Segoe UI'
            cor_fonte = RGBColor(115, 111, 111)
            texto = "#" + texto

        case "ibov":
            tamanho_fonte = Pt(12)
            nome_fonte = 'Segoe UI Semilight'
            cor_fonte = RGBColor(115, 111, 111)
            texto = texto + " pts"

        case "dolar":
            tamanho_fonte = Pt(20)
            nome_fonte = 'Segoe UI'
            cor_fonte = RGBColor(127, 127, 127)
            texto = "R$ " + texto + "*"

        case "data":
            tamanho_fonte = Pt(16)
            nome_fonte = 'Segoe UI Semilight'
            cor_fonte = RGBColor(115, 111, 111)

    paragrafo.alignment = PP_ALIGN.LEFT
    paragrafo.font.bold = bold
    paragrafo.font.name = nome_fonte
    paragrafo.font.size = tamanho_fonte
    paragrafo.font.color.rgb = cor_fonte
    paragrafo.text = texto

    shape.left = Pt(posicao_left)
    shape.top = Pt(posicao_top)

def inserir_negativo_positivo(
    slide: Slide,
    variacao: float,
    posicao_left: Pt,
    posicao_top: Pt,
    largura: Pt,
    altura: Pt,
) -> None:
    """Insere a imagem de negativo ou positivo no PDF do Fechamento de Mercado.

    ### Argumentos:
    - slide (Slide): Slide onde será adicionada a imagem
    - variacao (float): Valor da variação
    - posicao_left (Pt): Posição horizontal da imagem
    - posicao_top (Pt): Posição vertical da imagem
    - largura (Pt): Largura da imagem
    - altura (Pt): Altura da imagem"""

    sinal_negativo = str(BASE_DIR / 'source/assets/sinal_negativo.png')
    sinal_positivo = str(BASE_DIR / 'source/assets/sinal_positivo.png')

    picture = slide.shapes.add_picture(
        sinal_positivo if variacao >= 0 else sinal_negativo,
        Pt(posicao_left),
        Pt(posicao_top),
        Pt(largura),
        Pt(altura)
    )

    picture.left = Pt(posicao_left)
    picture.top = Pt(posicao_top)

def atualizar_data(
        shape: SlideShapes,
        data_atual: datetime.datetime,
        posicao_left: Pt,
        posicao_top: Pt,
) -> None:
    """Atualiza a data dentro do PDF do Fechamento do Mercado.

    ### Argumentos:
    - shape (SlidesShape): Shape onde será adicionado o texto
    - data_atual (datetime.datetime): Data atual
    - posicao_left (Pt): Posição horizontal do texto
    - posicao_top (Pt): Posição vertical do texto"""

    text_frame = shape.text_frame
    text_frame.paragraphs[0].font.size = Pt(16)
    text_frame.paragraphs[0].font.name = 'Segoe UI Semilight'
    text_frame.paragraphs[0].font.color.rgb = RGBColor(115, 111, 111)
    text_frame.paragraphs[0].text = data_atual.strftime('%d de %B de %Y')

    shape.left = Pt(posicao_left)
    shape.top = Pt(posicao_top)
